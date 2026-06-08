"""
ÉTAPE 2 — Extraction OPTIMISÉE de texte de tous les fichiers
============================================================
Architecture "fire-all-then-collect" pour Textract :
  Phase 1 : Extraction directe (Word, Excel, Email, Texte, PDF natifs) — séquentiel rapide
  Phase 2 : Batch upload S3 des fichiers nécessitant OCR — parallèle (20 threads)
  Phase 3 : Lancement en rafale de tous les jobs Textract async — rate-limited
  Phase 4 : Polling parallèle et collecte des résultats
  Phase 5 : Nettoyage S3 + checkpoint

Gain attendu : 5-10x plus rapide que la version séquentielle
Usage :
  python 02_extraction_optimized.py --copro 5033    # Mode per-copro (recommandé)
  python 02_extraction_optimized.py                  # Mode legacy (chemins hardcodés)
"""
import os
import sys
import json
import time
import shutil
import argparse
import boto3
import logging
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from botocore.config import Config
from tqdm import tqdm

from pipeline_config import paths_for

# ── Dépendances (installation auto si manquantes) ──
try:
    import fitz  # PyMuPDF
except ImportError:
    os.system("pip install PyMuPDF")
    import fitz
try:
    from docx import Document as DocxDocument
except ImportError:
    os.system("pip install python-docx")
    from docx import Document as DocxDocument
try:
    import openpyxl
except ImportError:
    os.system("pip install openpyxl")
    import openpyxl
try:
    import extract_msg
except ImportError:
    os.system("pip install extract-msg")
    import extract_msg
try:
    from pptx import Presentation
except ImportError:
    os.system("pip install python-pptx")
    from pptx import Presentation


# =====================================================
# CONFIGURATION
# =====================================================
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

_parser = argparse.ArgumentParser(description="Extraction de texte d'une copropriété.")
_parser.add_argument("--copro", help="Code NCG de la copropriété (ex: 5033). Si absent, mode legacy.")
_args, _ = _parser.parse_known_args()

if _args.copro:
    _paths = paths_for(_args.copro)
    # FILTERED_DIR / OUTPUT_DIR pointent sur la base (Archives_Filtrees / Archives_Extraites)
    # pour que rel_path reste au format legacy "{copro_folder}/.../file.pdf" et que les
    # champs copropriete/dossier_parent/source_file restent compatibles avec le mode legacy.
    # On restreint juste le walk au sous-dossier copro via _WALK_DIR.
    FILTERED_DIR = str(_paths["filtered"].parent)   # = Archives_Filtrees
    OUTPUT_DIR   = str(_paths["extracted"].parent)  # = Archives_Extraites
    _WALK_DIR    = str(_paths["filtered"])          # = Archives_Filtrees/{copro}
    _paths["per_copro"].mkdir(parents=True, exist_ok=True)
    CHECKPOINT_FILE = str(_paths["extraction_checkpoint"])
    _LOG_PATH = str(_paths["extraction_log"])
    print(f"📌 Mode per-copro : {_args.copro} ({_paths['folder_name']})")
else:
    FILTERED_DIR = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Filtrees"
    OUTPUT_DIR   = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Extraites"
    _WALK_DIR    = FILTERED_DIR
    CHECKPOINT_FILE = os.path.join(SCRIPT_DIR, "extraction_checkpoint.json")
    _LOG_PATH = "extraction.log"

S3_BUCKET = "smarterplan-rag-prototype"
S3_TEXTRACT_PREFIX = "textract_temp/"
AWS_REGION = "eu-west-1"

# ── Paramètres de performance ──
MAX_CONCURRENT_UPLOADS = 20       # Threads pour upload S3
JOB_LAUNCH_RATE = 5               # Jobs lancés par seconde (quota Textract défaut)
POLL_INTERVAL = 3                 # Secondes entre chaque cycle de polling
POLL_BATCH_SIZE = 50              # Jobs vérifiés par cycle
MAX_RETRIES = 3                   # Retries sur throttling

# ── AWS clients avec retry adaptatif ──
boto_config = Config(
    retries={"max_attempts": 10, "mode": "adaptive"},
    max_pool_connections=50
)
textract = boto3.client("textract", region_name=AWS_REGION, config=boto_config)
s3 = boto3.client("s3", region_name=AWS_REGION, config=boto_config)

# ── Logging ──
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(_LOG_PATH, encoding="utf-8")
    ]
)
log = logging.getLogger(__name__)

# ── Stats ──
stats = {
    "pdf_natif": 0, "pdf_ocr": 0, "word": 0, "excel": 0,
    "email": 0, "texte": 0, "pptx": 0, "image_ocr": 0,
    "erreurs": 0, "vides": 0, "skipped_checkpoint": 0
}


# =====================================================
# CHECKPOINT : reprise après interruption
# =====================================================
def sig_for(rel_path):
    """Signature de contenu bon marche (taille:mtime_ns) du fichier source de rel_path.
    Change si le doc est modifie (meme nom) -> declenche la re-extraction (U du CRUD)."""
    try:
        st = os.stat(os.path.join(FILTERED_DIR, rel_path))
        return f"{st.st_size}:{st.st_mtime_ns}"
    except OSError:
        return None

def load_checkpoint():
    """Retourne {rel_path: signature}. signature=None => extraction d'avant le suivi
    par signature : baseline adoptee au prochain passage SANS re-extraire."""
    sigs = {}
    # 1. Charger le checkpoint (nouveau format {sigs:{...}} ou ancien {completed:[...]})
    if os.path.exists(CHECKPOINT_FILE):
        try:
            with open(CHECKPOINT_FILE, "r") as f:
                data = json.load(f)
            if isinstance(data.get("sigs"), dict):
                sigs.update(data["sigs"])               # nouveau format {rel_path: sig}
            for rp in data.get("completed", []):         # ancien format (liste) -> sig inconnue
                sigs.setdefault(rp, None)
        except Exception:
            pass
    # 2. Scanner l'output pour ce qui est VRAIMENT fait (JSON present, sig inconnue)
    if os.path.exists(OUTPUT_DIR):
        for root, dirs, files in os.walk(OUTPUT_DIR):
            for file in files:
                if file.endswith(".json"):
                    rel_json_path = os.path.relpath(os.path.join(root, file), OUTPUT_DIR)
                    rel_source_path = rel_json_path[:-5] # enlever '.json'
                    if rel_source_path != "extraction_checkpoint":
                        sigs.setdefault(rel_source_path, None)
    return sigs

def save_checkpoint(sigs):
    with open(CHECKPOINT_FILE, "w") as f:
        json.dump({"sigs": sigs}, f)

completed_files = load_checkpoint()
if completed_files:
    log.info(f"♻️  Reprise : {len(completed_files)} fichiers déjà traités physiquement/dans le checkpoint")


# =====================================================
# EXTRACTEURS DIRECTS (non-Textract)
# =====================================================
def extract_pdf_native(filepath):
    """
    Retourne (texte, is_native). is_native=False → nécessite Textract.
    
    Détection robuste en 3 critères :
      1. Moyenne globale : >300 chars/page (un vrai PDF natif FR fait 1500-3000)
      2. Couverture : >80% des pages doivent avoir du texte substantiel (>100 chars)
         → détecte les PDFs mixtes (pages natives + pages scannées)
      3. Fallback : si le PDF a ≤2 pages, on utilise uniquement la moyenne
    
    Seuil 300 chars/page choisi car :
      - Un OCR embarqué de mauvaise qualité produit ~50-150 chars/page
      - Un PDF natif français avec du contenu juridique fait >1500 chars/page
      - Marge de sécurité pour les pages de garde, sommaires courts, etc.
    """
    NATIVE_THRESHOLD_CHARS_PER_PAGE = 300
    MIN_PAGE_CHARS = 100            # Seuil pour qu'une page soit "substantielle"
    MIN_COVERAGE_RATIO = 0.80       # 80% des pages doivent avoir du texte

    try:
        doc = fitz.open(filepath)
        page_count = doc.page_count

        if page_count == 0:
            doc.close()
            return "", False

        page_texts = []
        full_text = ""
        for page in doc:
            pt = page.get_text()
            page_texts.append(pt)
            full_text += pt
        doc.close()

        total_chars = len(full_text.strip())
        avg_chars = total_chars / page_count

        # Critère 1 : moyenne globale
        if avg_chars < NATIVE_THRESHOLD_CHARS_PER_PAGE:
            log.debug(f"PDF scanné (avg {avg_chars:.0f} chars/page < {NATIVE_THRESHOLD_CHARS_PER_PAGE}): {filepath}")
            return "", False

        # Critère 2 : couverture page par page (sauf docs très courts)
        if page_count > 2:
            pages_with_text = sum(1 for pt in page_texts if len(pt.strip()) >= MIN_PAGE_CHARS)
            coverage = pages_with_text / page_count

            if coverage < MIN_COVERAGE_RATIO:
                log.info(f"PDF mixte détecté ({coverage:.0%} couverture, {pages_with_text}/{page_count} pages) → OCR: {os.path.basename(filepath)}")
                return "", False

        return full_text.strip(), True

    except Exception as e:
        log.debug(f"Erreur lecture PDF {filepath}: {e}")
        return "", False

def extract_docx(filepath):
    try:
        doc = DocxDocument(filepath)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception:
        try:
            with open(filepath, "rb") as f:
                raw = f.read()
            text = raw.decode("utf-8", errors="ignore")
            return "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        except Exception:
            return ""

def extract_excel(filepath):
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            texts.append(f"--- Feuille: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join([str(c) for c in row if c is not None])
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_email_msg(filepath):
    try:
        msg = extract_msg.Message(filepath)
        parts = []
        if msg.subject: parts.append(f"Objet: {msg.subject}")
        if msg.sender: parts.append(f"De: {msg.sender}")
        if msg.date: parts.append(f"Date: {msg.date}")
        if msg.body: parts.append(f"\n{msg.body}")
        return "\n".join(parts)
    except Exception:
        return ""

def extract_pptx(filepath):
    try:
        prs = Presentation(filepath)
        texts = []
        for i, slide in enumerate(prs.slides):
            texts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_file(filepath):
    for enc in ["utf-8", "latin-1", "cp1252"]:
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            continue
    return ""


DIRECT_EXTRACTORS = {
    "word": extract_docx,
    "excel": extract_excel,
    "email": extract_email_msg,
    "texte": extract_text_file,
    "pptx": extract_pptx,
}

DIRECT_EXT_MAP = {
    ".doc": "word", ".docx": "word",
    ".xls": "excel", ".xlsx": "excel",
    ".msg": "email", ".eml": "email",
    ".txt": "texte", ".rtf": "texte", ".csv": "texte",
    ".ppt": "pptx", ".pptx": "pptx",
}

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp"}


# =====================================================
# HELPERS
# =====================================================
def save_extracted(rel_path, ext, text, output_dir):
    """Sauvegarde un résultat d'extraction en JSON."""
    if not text or len(text.strip()) < 20:
        stats["vides"] += 1
        return False

    output_data = {
        "source_file": rel_path,
        "source_extension": ext,
        "copropriete": rel_path.split(os.sep)[0],
        "dossier_parent": os.path.dirname(rel_path),
        "nom_fichier": os.path.basename(rel_path),
        "texte": text,
        "nb_caracteres": len(text)
    }
    output_path = os.path.join(output_dir, rel_path + ".json")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)
    return True


# =====================================================
# PHASE 0 : TRIAGE
# =====================================================
def triage_files():
    """Sépare fichiers en extraction directe vs OCR Textract."""
    direct_files = []   # (filepath, rel_path, ext, extractor_name)
    ocr_files = []      # (filepath, rel_path, ext)

    log.info("Phase 0: Triage des fichiers...")
    for root, dirs, filenames in os.walk(_WALK_DIR):
        for fname in filenames:
            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, FILTERED_DIR)
            ext = Path(fname).suffix.lower()

            # Skip sensible au CONTENU (U du CRUD) : on saute seulement si le JSON
            # existe ET la signature du fichier source est inchangee. Signature
            # differente = modification sur place -> re-extraction.
            output_json_path = os.path.join(OUTPUT_DIR, rel_path + ".json")
            json_exists = os.path.exists(output_json_path)
            cur_sig = sig_for(rel_path)
            if json_exists and rel_path in completed_files:
                prev_sig = completed_files[rel_path]
                if prev_sig is None:
                    # Migration : JSON present, sig inconnue -> adopter le contenu actuel
                    # comme baseline, sans re-extraire (pas de re-OCR du corpus existant).
                    completed_files[rel_path] = cur_sig
                    stats["skipped_checkpoint"] += 1
                    continue
                if prev_sig == cur_sig:
                    stats["skipped_checkpoint"] += 1
                    continue
                # Signature differente -> MODIFICATION -> re-extraire (purge l'ancien JSON).
                log.info(f"Modif detectee, re-extraction : {rel_path}")
                try:
                    os.remove(output_json_path)
                except OSError:
                    pass
                completed_files.pop(rel_path, None)
            elif json_exists:
                # JSON present mais inconnu du checkpoint -> adopter baseline, skip.
                completed_files[rel_path] = cur_sig
                stats["skipped_checkpoint"] += 1
                continue
            elif rel_path in completed_files:
                # JSON manquant mais dans le checkpoint -> re-traitement force.
                log.info(f"Fichier cible manquant, re-traitement force : {rel_path}")
                completed_files.pop(rel_path, None)

            if ext in DIRECT_EXT_MAP:
                direct_files.append((filepath, rel_path, ext, DIRECT_EXT_MAP[ext]))
            elif ext == ".pdf":
                text, is_native = extract_pdf_native(filepath)
                if is_native:
                    # Stocker le texte déjà extrait pour éviter double lecture
                    direct_files.append((filepath, rel_path, ext, "pdf_natif"))
                else:
                    ocr_files.append((filepath, rel_path, ext))
            elif ext in IMAGE_EXTS:
                ocr_files.append((filepath, rel_path, ext))

    return direct_files, ocr_files


# =====================================================
# PHASE 1 : EXTRACTION DIRECTE
# =====================================================
def run_direct_extraction(direct_files):
    log.info(f"\nPhase 1: Extraction directe de {len(direct_files)} fichiers...")
    for filepath, rel_path, ext, etype in tqdm(direct_files, desc="Extraction directe"):
        if etype == "pdf_natif":
            text, _ = extract_pdf_native(filepath)
        else:
            text = DIRECT_EXTRACTORS[etype](filepath)

        if save_extracted(rel_path, ext, text, OUTPUT_DIR):
            stats[etype] += 1
            completed_files[rel_path] = sig_for(rel_path)
        else:
            if not text or len(text.strip()) < 20:
                pass  # déjà compté dans save_extracted
            else:
                stats["erreurs"] += 1

    save_checkpoint(completed_files)
    log.info(f"✅ Phase 1 terminée")


# =====================================================
# PHASE 2 : UPLOAD S3 PARALLÈLE
# =====================================================
def upload_one(filepath, s3_key):
    try:
        s3.upload_file(filepath, S3_BUCKET, s3_key)
        return s3_key
    except Exception as e:
        log.error(f"Upload échoué {os.path.basename(filepath)}: {e}")
        return None

def batch_upload_s3(ocr_files):
    log.info(f"\nPhase 2: Upload S3 de {len(ocr_files)} fichiers ({MAX_CONCURRENT_UPLOADS} threads)...")
    s3_map = {}  # rel_path → (s3_key, ext)

    with ThreadPoolExecutor(max_workers=MAX_CONCURRENT_UPLOADS) as executor:
        futures = {}
        for filepath, rel_path, ext in ocr_files:
            s3_key = S3_TEXTRACT_PREFIX + rel_path.replace(os.sep, "/")
            fut = executor.submit(upload_one, filepath, s3_key)
            futures[fut] = (rel_path, s3_key, ext)

        for fut in tqdm(as_completed(futures), total=len(futures), desc="Upload S3"):
            rel_path, s3_key, ext = futures[fut]
            if fut.result():
                s3_map[rel_path] = (s3_key, ext)

    log.info(f"✅ {len(s3_map)}/{len(ocr_files)} fichiers sur S3")
    return s3_map


# =====================================================
# PHASE 3 : LANCEMENT RAFALE TEXTRACT
# =====================================================
def _reconstruct_from_word_blocks(textract_response):
    """Reconstruit le texte depuis les WORD blocks d'une réponse sync Textract."""
    words = []
    for block in textract_response.get("Blocks", []):
        if block["BlockType"] == "WORD":
            bbox = block.get("Geometry", {}).get("BoundingBox", {})
            words.append({
                "text": block["Text"],
                "top": bbox.get("Top", 0),
                "left": bbox.get("Left", 0),
                "height": bbox.get("Height", 0),
            })

    if not words:
        return ""

    words.sort(key=lambda w: (w["top"], w["left"]))

    lines = []
    current_line = [words[0]]
    avg_height = words[0]["height"] or 0.01

    for w in words[1:]:
        tolerance = avg_height * 0.4
        if abs(w["top"] - current_line[-1]["top"]) <= tolerance:
            current_line.append(w)
        else:
            lines.append(current_line)
            current_line = [w]
        if w["height"] > 0:
            avg_height = avg_height * 0.8 + w["height"] * 0.2

    lines.append(current_line)

    text_parts = []
    prev_bottom = 0
    for line_words in lines:
        line_words.sort(key=lambda w: w["left"])
        line_text = " ".join(w["text"] for w in line_words)

        line_top = line_words[0]["top"]
        line_h = max(w["height"] for w in line_words)
        if prev_bottom > 0 and (line_top - prev_bottom) > line_h * 1.2:
            text_parts.append("")

        text_parts.append(line_text)
        prev_bottom = line_top + line_h

    return "\n".join(text_parts)


def start_textract_job(s3_key, ext):
    """Lance un job. Retourne ('ASYNC', job_id) ou ('SYNC', text) ou None."""
    for attempt in range(MAX_RETRIES):
        try:
            if ext == ".pdf":
                resp = textract.start_document_text_detection(
                    DocumentLocation={"S3Object": {"Bucket": S3_BUCKET, "Name": s3_key}}
                )
                return ("ASYNC", resp["JobId"])
            else:
                # Images ≤10MB : API synchrone (plus rapide, pas de polling)
                obj = s3.get_object(Bucket=S3_BUCKET, Key=s3_key)
                img_bytes = obj["Body"].read()

                if len(img_bytes) > 10_000_000:
                    resp = textract.start_document_text_detection(
                        DocumentLocation={"S3Object": {"Bucket": S3_BUCKET, "Name": s3_key}}
                    )
                    return ("ASYNC", resp["JobId"])

                resp = textract.detect_document_text(Document={"Bytes": img_bytes})
                text = _reconstruct_from_word_blocks(resp)
                return ("SYNC", text)

        except textract.exceptions.ProvisionedThroughputExceededException:
            time.sleep(2 ** attempt)
        except Exception as e:
            log.error(f"Textract start error {s3_key}: {e}")
            return None
    return None

def launch_all_jobs(s3_map):
    log.info(f"\nPhase 3: Lancement de {len(s3_map)} jobs Textract (rate: {JOB_LAUNCH_RATE}/s)...")
    async_jobs = {}      # job_id → (rel_path, s3_key, ext)
    sync_results = {}    # rel_path → (text, ext)
    failed = []

    for i, (rel_path, (s3_key, ext)) in enumerate(
        tqdm(s3_map.items(), desc="Lancement Textract")
    ):
        result = start_textract_job(s3_key, ext)

        if result is None:
            failed.append(rel_path)
        elif result[0] == "ASYNC":
            async_jobs[result[1]] = (rel_path, s3_key, ext)
        elif result[0] == "SYNC":
            sync_results[rel_path] = (result[1], ext)

        # Rate limiting
        if (i + 1) % JOB_LAUNCH_RATE == 0:
            time.sleep(1.0)

    log.info(f"✅ {len(async_jobs)} async + {len(sync_results)} sync + {len(failed)} échecs")
    return async_jobs, sync_results, failed


# =====================================================
# PHASE 4 : POLLING DES RÉSULTATS ASYNC
# =====================================================
def collect_result(job_id):
    """
    Récupère les résultats d'un job Textract et reconstruit le texte
    en utilisant les WORD blocks + coordonnées géométriques.
    
    Les LINE blocks collent parfois les mots sur les scans de mauvaise
    qualité. Les WORD blocks sont toujours correctement segmentés —
    on reconstruit les lignes via les coordonnées BoundingBox.
    """
    words = []
    kwargs = {"JobId": job_id}

    while True:
        result = textract.get_document_text_detection(**kwargs)
        for block in result.get("Blocks", []):
            if block["BlockType"] == "WORD":
                bbox = block.get("Geometry", {}).get("BoundingBox", {})
                words.append({
                    "text": block["Text"],
                    "page": block.get("Page", 1),
                    "top": bbox.get("Top", 0),
                    "left": bbox.get("Left", 0),
                    "height": bbox.get("Height", 0),
                })
        next_token = result.get("NextToken")
        if not next_token:
             break
        kwargs["NextToken"] = next_token

    if not words:
        return ""

    # Grouper par page
    pages = {}
    for w in words:
        pages.setdefault(w["page"], []).append(w)

    full_text_parts = []

    for page_num in sorted(pages.keys()):
        page_words = pages[page_num]
        page_words.sort(key=lambda w: (w["top"], w["left"]))

        # Grouper en lignes : mots avec un "top" similaire (± tolérance)
        lines = []
        current_line = [page_words[0]]
        avg_height = page_words[0]["height"] or 0.01

        for w in page_words[1:]:
            tolerance = avg_height * 0.4
            if abs(w["top"] - current_line[-1]["top"]) <= tolerance:
                current_line.append(w)
            else:
                lines.append(current_line)
                current_line = [w]
            if w["height"] > 0:
                avg_height = avg_height * 0.8 + w["height"] * 0.2

        lines.append(current_line)

        # Reconstruire le texte ligne par ligne
        prev_line_bottom = 0
        for line_words in lines:
            line_words.sort(key=lambda w: w["left"])
            line_text = " ".join(w["text"] for w in line_words)

            # Détecter les sauts de paragraphe (gap vertical important)
            line_top = line_words[0]["top"]
            line_height = max(w["height"] for w in line_words) if line_words else 0.01
            gap = line_top - prev_line_bottom

            if prev_line_bottom > 0 and gap > line_height * 1.2:
                full_text_parts.append("")  # ligne vide = saut de paragraphe

            full_text_parts.append(line_text)
            prev_line_bottom = line_top + line_height

    return "\n".join(full_text_parts)

def poll_all_jobs(async_jobs):
    log.info(f"\nPhase 4: Polling de {len(async_jobs)} jobs...")
    pending = dict(async_jobs)
    results = {}
    failed = []
    cycle = 0

    while pending:
        cycle += 1
        done_this_round = []
        job_ids = list(pending.keys())[:POLL_BATCH_SIZE]

        for job_id in job_ids:
            try:
                resp = textract.get_document_text_detection(JobId=job_id, MaxResults=1)
                status = resp["JobStatus"]

                if status == "SUCCEEDED":
                    rel_path, s3_key, ext = pending[job_id]
                    text = collect_result(job_id)
                    results[rel_path] = (text, ext)
                    done_this_round.append(job_id)
                elif status == "FAILED":
                    rel_path, _, _ = pending[job_id]
                    log.warning(f"FAILED: {rel_path}")
                    failed.append(rel_path)
                    done_this_round.append(job_id)

            except Exception as e:
                log.warning(f"Poll error: {e}")

        for jid in done_this_round:
            del pending[jid]

        if pending:
            done_n = len(results) + len(failed)
            total_n = done_n + len(pending)
            log.info(f"  Cycle {cycle}: {done_n}/{total_n} ({done_n/total_n*100:.0f}%) — {len(pending)} restants")
            time.sleep(POLL_INTERVAL)

    log.info(f"✅ Polling terminé: {len(results)} OK, {len(failed)} KO")
    return results, failed


# =====================================================
# PHASE 5 : NETTOYAGE S3
# =====================================================
def cleanup_s3(s3_map):
    log.info("\nPhase 5: Nettoyage S3...")
    keys = [s3_key for s3_key, _ in s3_map.values()]
    for i in range(0, len(keys), 1000):
        batch = keys[i:i+1000]
        try:
            s3.delete_objects(
                Bucket=S3_BUCKET,
                Delete={"Objects": [{"Key": k} for k in batch]}
            )
        except Exception as e:
            log.warning(f"Nettoyage S3 partiel: {e}")
    log.info(f"✅ {len(keys)} fichiers temp supprimés")


# =====================================================
# MAIN
# =====================================================
def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    print("=" * 60)
    print("  EXTRACTION OPTIMISÉE — PIPELINE PARALLÈLE")
    print("=" * 60)

    # ── Triage ──
    direct_files, ocr_files = triage_files()
    log.info(f"  {len(direct_files)} directs | {len(ocr_files)} OCR | {stats['skipped_checkpoint']} checkpoint")

    # ── Phase 1 : Extraction directe ──
    run_direct_extraction(direct_files)

    if not ocr_files:
        log.info("Aucun fichier OCR. Terminé !")
        print_report()
        return

    # ── Phase 2 : Upload S3 ──
    s3_map = batch_upload_s3(ocr_files)

    # ── Phase 3 : Lancement Textract ──
    async_jobs, sync_results, launch_failed = launch_all_jobs(s3_map)

    # Sauvegarder résultats synchrones immédiatement
    for rel_path, (text, ext) in sync_results.items():
        if save_extracted(rel_path, ext, text, OUTPUT_DIR):
            stats["image_ocr"] += 1
            completed_files[rel_path] = sig_for(rel_path)
    save_checkpoint(completed_files)

    # ── Phase 4 : Polling async ──
    if async_jobs:
        ocr_results, poll_failed = poll_all_jobs(async_jobs)
        for rel_path, (text, ext) in ocr_results.items():
            if save_extracted(rel_path, ext, text, OUTPUT_DIR):
                if ext == ".pdf":
                    stats["pdf_ocr"] += 1
                else:
                    stats["image_ocr"] += 1
                completed_files[rel_path] = sig_for(rel_path)
        stats["erreurs"] += len(poll_failed) + len(launch_failed)
        save_checkpoint(completed_files)

    # ── Phase 5 : Nettoyage ──
    cleanup_s3(s3_map)

    # ── Rapport ──
    print_report()

    if stats["erreurs"] == 0 and os.path.exists(CHECKPOINT_FILE):
        os.remove(CHECKPOINT_FILE)
    elif stats["erreurs"] > 0:
        print(f"\n  ⚠️  {stats['erreurs']} erreurs. Relance le script pour retenter les fichiers échoués.")


def print_report():
    print("\n" + "=" * 60)
    print("  RAPPORT D'EXTRACTION")
    print("=" * 60)
    print(f"  PDF natifs (texte direct) : {stats['pdf_natif']}")
    print(f"  PDF scannés (Textract)    : {stats['pdf_ocr']}")
    print(f"  Word                      : {stats['word']}")
    print(f"  Excel                     : {stats['excel']}")
    print(f"  Emails                    : {stats['email']}")
    print(f"  Texte/CSV                 : {stats['texte']}")
    print(f"  PowerPoint                : {stats['pptx']}")
    print(f"  Images OCR (plans)        : {stats['image_ocr']}")
    print(f"  Fichiers vides/courts     : {stats['vides']}")
    print(f"  Erreurs                   : {stats['erreurs']}")
    t = stats['pdf_ocr'] + stats['image_ocr']
    print(f"\n  TOTAL Textract            : {t} fichiers")
    print(f"  Coût Textract estimé      : ~${t * 0.0015:.2f}")
    print(f"\n  📁 Résultats : {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
