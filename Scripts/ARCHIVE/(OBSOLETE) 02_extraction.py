"""
ÉTAPE 2 — Extraction de texte de tous les fichiers
- PDF scannés → Textract
- PDF natifs → PyMuPDF (fitz)
- Word/Excel/Email/Texte → extraction directe
Lance : python 02_extraction.py
"""
import os
import sys
import json
import time
import boto3
from pathlib import Path

# Pour les différents formats
try:
    import fitz  # PyMuPDF - meilleur que PyPDF2
except ImportError:
    print("Installation de PyMuPDF...")
    os.system("pip install PyMuPDF")
    import fitz

try:
    from docx import Document as DocxDocument
except ImportError:
    os.system("pip install python-docx")
    from docx import Document as DocxDocument

try:
    import openpyxl
except ImportError:
    os.system("pip install openpyxl")
    import openpyxl

try:
    import extract_msg
except ImportError:
    os.system("pip install extract-msg")
    import extract_msg

try:
    from pptx import Presentation
except ImportError:
    os.system("pip install python-pptx")
    from pptx import Presentation

try:
    from tqdm import tqdm
except ImportError:
    os.system("pip install tqdm")
    from tqdm import tqdm

# =====================================================
# CONFIGURATION
# =====================================================
FILTERED_DIR = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Filtrees"  # ← MODIFIER
OUTPUT_DIR = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Extraites"    # ← MODIFIER
S3_BUCKET = "smarterplan-rag-prototype"
AWS_REGION = "eu-west-1"

# Textract client
textract = boto3.client("textract", region_name=AWS_REGION)
s3 = boto3.client("s3", region_name=AWS_REGION)

# Compteurs et logs
stats = {"pdf_natif": 0, "pdf_ocr": 0, "word": 0, "excel": 0, 
         "email": 0, "texte": 0, "pptx": 0, "image_ocr": 0, "erreurs": 0, "vides": 0}
failed_files = [] # Liste pour garder trace des fichiers en erreur

# =====================================================
# Fonctions d'extraction par type
# =====================================================

def extract_pdf(filepath):
    """Tente extraction texte natif. Si échec ou vide → Textract OCR."""
    try:
        doc = fitz.open(filepath)
        text = ""
        page_count = doc.page_count
        for page in doc:
            text += page.get_text()
        doc.close()
        
        # Si le PDF contient du texte sélectionnable (>50 chars par page en moyenne)
        if len(text.strip()) > 50 * page_count:
            stats["pdf_natif"] += 1
            return text.strip()
    except Exception as e:
        print(f"  ⚠️ Erreur lecture PDF natif {os.path.basename(filepath)}: {e}")
    
    # Sinon c'est un scan → Textract
    return extract_via_textract(filepath)

def extract_via_textract(filepath):
    """Envoie un fichier à Textract pour OCR."""
    filename = os.path.basename(filepath)
    s3_key = f"textract_temp/{filename}"
    
    # Upload temporaire vers S3 (Textract ne lit pas en local)
    try:
        s3.upload_file(filepath, S3_BUCKET, s3_key)
    except Exception as e:
        print(f"  ⚠️ Erreur upload S3 pour {filename}: {e}")
        stats["erreurs"] += 1
        failed_files.append((filepath, str(e)))
        return ""
    
    try:
        # Pour les documents > 1 page, utiliser l'API asynchrone
        ext = Path(filepath).suffix.lower()
        
        if ext == ".pdf":
            # API asynchrone pour PDF multi-pages
            response = textract.start_document_text_detection(
                DocumentLocation={"S3Object": {"Bucket": S3_BUCKET, "Name": s3_key}},
                OutputConfig={"S3Bucket": S3_BUCKET, "S3Prefix": "textract_output/"}
            )
            job_id = response["JobId"]
            
            # Attendre la fin du job
            while True:
                result = textract.get_document_text_detection(JobId=job_id)
                status = result["JobStatus"]
                if status == "SUCCEEDED":
                    break
                elif status == "FAILED":
                    stats["erreurs"] += 1
                    return ""
                time.sleep(2)
            
            # Récupérer toutes les pages de résultats
            text_blocks = []
            next_token = None
            while True:
                if next_token:
                    result = textract.get_document_text_detection(JobId=job_id, NextToken=next_token)
                else:
                    result = textract.get_document_text_detection(JobId=job_id)

                for block in result.get("Blocks", []):
                    if block["BlockType"] == "LINE":
                        text_blocks.append(block["Text"])
                
                next_token = result.get("NextToken")
                if not next_token:
                    break
            
            stats["pdf_ocr"] += 1
            return "\n".join(text_blocks)
        
        else:
            # API synchrone pour images (1 page)
            with open(filepath, "rb") as f:
                img_bytes = f.read()
            
            response = textract.detect_document_text(
                Document={"Bytes": img_bytes}
            )
            
            text_blocks = [
                block["Text"] for block in response.get("Blocks", [])
                if block["BlockType"] == "LINE"
            ]
            
            stats["image_ocr"] += 1
            return "\n".join(text_blocks)
    
    except Exception as e:
        print(f"  ⚠️ Erreur Textract pour {filename}: {e}")
        stats["erreurs"] += 1
        failed_files.append((filepath, str(e)))
        return ""
    
    finally:
        # Nettoyer le fichier temporaire S3
        try:
            s3.delete_object(Bucket=S3_BUCKET, Key=s3_key)
        except:
            pass

def extract_docx(filepath):
    """Extraction texte d'un fichier Word."""
    try:
        doc = DocxDocument(filepath)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        stats["word"] += 1
        return text
    except Exception as e:
        # Ancien format .doc → tenter lecture brute
        try:
            with open(filepath, "rb") as f:
                raw = f.read()
            # Extraction basique du texte depuis le binaire .doc
            text = raw.decode("utf-8", errors="ignore")
            # Nettoyer les caractères non-texte
            text = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
            stats["word"] += 1
            return text
        except:
            stats["erreurs"] += 1
            return ""

def extract_excel(filepath):
    """Extraction texte d'un fichier Excel — toutes les feuilles."""
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            texts.append(f"--- Feuille: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    texts.append(row_text)
        stats["excel"] += 1
        return "\n".join(texts)
    except Exception as e:
        stats["erreurs"] += 1
        return ""

def extract_email(filepath):
    """Extraction texte d'un fichier .msg."""
    try:
        msg = extract_msg.Message(filepath)
        parts = []
        if msg.subject:
            parts.append(f"Objet: {msg.subject}")
        if msg.sender:
            parts.append(f"De: {msg.sender}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        if msg.body:
            parts.append(f"\n{msg.body}")
        stats["email"] += 1
        return "\n".join(parts)
    except Exception as e:
        stats["erreurs"] += 1
        return ""

def extract_pptx(filepath):
    """Extraction texte d'un PowerPoint."""
    try:
        prs = Presentation(filepath)
        texts = []
        for i, slide in enumerate(prs.slides):
            texts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        stats["pptx"] += 1
        return "\n".join(texts)
    except:
        stats["erreurs"] += 1
        return ""

def extract_text_file(filepath):
    """Lecture d'un fichier texte."""
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(filepath, "r", encoding=enc) as f:
                stats["texte"] += 1
                return f.read()
        except:
            continue
    stats["erreurs"] += 1
    return ""

# =====================================================
# Mapping extension → fonction d'extraction
# =====================================================
EXTRACTORS = {
    ".pdf": extract_pdf,
    ".doc": extract_docx,
    ".docx": extract_docx,
    ".xls": extract_excel,
    ".xlsx": extract_excel,
    ".msg": extract_email,
    ".eml": extract_email,
    ".txt": extract_text_file,
    ".rtf": extract_text_file,
    ".csv": extract_text_file,
    ".ppt": extract_pptx,
    ".pptx": extract_pptx,
    ".jpg": extract_via_textract,
    ".jpeg": extract_via_textract,
    ".png": extract_via_textract,
    ".tif": extract_via_textract,
    ".tiff": extract_via_textract,
    ".bmp": extract_via_textract,
}

# =====================================================
# Exécution
# =====================================================
# Nettoyage des anciennes extractions pour repartir de zéro
if os.path.exists(OUTPUT_DIR):
    print(f"Nettoyage du dossier de sortie : {OUTPUT_DIR}")
    import shutil
    shutil.rmtree(OUTPUT_DIR)

print("=" * 50)
print("EXTRACTION DE TEXTE — TOUS FORMATS")
print("=" * 50)

# S'assurer que le dossier de sortie existe
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Lister tous les fichiers
all_files = []
if os.path.exists(FILTERED_DIR):
    for root, dirs, filenames in os.walk(FILTERED_DIR):
        for fname in filenames:
            filepath = os.path.join(root, fname)
            ext = Path(fname).suffix.lower()
            if ext in EXTRACTORS:
                all_files.append((filepath, ext))
else:
    print(f"❌ Le dossier {FILTERED_DIR} n'existe pas. Lance d'abord le script de filtrage.")
    sys.exit(1)

print(f"\n{len(all_files)} fichiers à traiter\n")

for filepath, ext in tqdm(all_files, desc="Extraction"):
    rel_path = os.path.relpath(filepath, FILTERED_DIR)
    
    # Extraire le texte
    extractor = EXTRACTORS[ext]
    text = extractor(filepath)
    
    if not text or len(text.strip()) < 20:
        stats["vides"] += 1
        continue
    
    # Sauvegarder le texte extrait en JSON avec métadonnées
    output_data = {
        "source_file": rel_path,
        "source_extension": ext,
        "copropriete": rel_path.split(os.sep)[0] or "RACINE",
        "dossier_parent": os.path.dirname(rel_path),
        "nom_fichier": os.path.basename(filepath),
        "texte": text,
        "nb_caracteres": len(text)
    }
    
    # Sauvegarder
    output_path = os.path.join(OUTPUT_DIR, rel_path + ".json")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)

# =====================================================
# Rapport final
# =====================================================
print("\n" + "=" * 50)
print("RAPPORT D'EXTRACTION")
print("=" * 50)
print(f"  PDF natifs (texte direct) : {stats['pdf_natif']}")
print(f"  PDF scannés (Textract)    : {stats['pdf_ocr']}")
print(f"  Word                      : {stats['word']}")
print(f"  Excel                     : {stats['excel']}")
print(f"  Emails                    : {stats['email']}")
print(f"  Texte/CSV                 : {stats['texte']}")
print(f"  PowerPoint                : {stats['pptx']}")
print(f"  Images OCR (plans)        : {stats['image_ocr']}")
print(f"  Fichiers vides/très courts: {stats['vides']}")
print(f"  Erreurs                   : {stats['erreurs']}")
print(f"\n  TOTAL Textract utilisé    : {stats['pdf_ocr'] + stats['image_ocr']} fichiers")
print(f"  Coût Textract estimé      : ~${(stats['pdf_ocr'] + stats['image_ocr']) * 0.0015:.2f}")
print(f"\n📁 Textes extraits dans : {OUTPUT_DIR}")

# Sauvegarder la liste des erreurs si besoin
if failed_files:
    ERROR_LOG = os.path.join(SCRIPT_DIR if 'SCRIPT_DIR' in locals() else ".", "erreurs_extraction.txt")
    with open(ERROR_LOG, "w", encoding="utf-8") as f:
        f.write("LISTE DES FICHIERS EN ERREUR\n")
        f.write("="*30 + "\n")
        for path, err in failed_files:
            f.write(f"- {path}\n  Erreur: {err}\n\n")
    print(f"📋 Liste détaillée des {len(failed_files)} erreurs dans : {ERROR_LOG}")
