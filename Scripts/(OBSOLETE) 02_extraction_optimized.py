"""
ÉTAPE 2 — Extraction OPTIMISÉE de texte de tous les fichiers
============================================================
Architecture "fire-all-then-collect" pour Textract :
  Phase 1 : Extraction directe (Word, Excel, Email, Texte, PDF natifs) — séquentiel rapide
  Phase 2 : Batch upload S3 des fichiers nécessitant OCR — parallèle (20 threads)
  Phase 3 : Lancement en rafale de tous les jobs Textract async — rate-limited
  Phase 4 : Polling parallèle et collecte des résultats
  Phase 5 : Nettoyage S3 + checkpoint

Gain attendu : 5-10x plus rapide que la version séquentielle
Lance : python 02_extraction_optimized.py
"""
import os
import sys
import json
import time
import shutil
import boto3
import logging
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from botocore.config import Config
from tqdm import tqdm

# ── Dépendances (installation auto si manquantes) ──
try:
    import fitz  # PyMuPDF
except ImportError:
    os.system("pip install PyMuPDF")
    import fitz
try:
    from docx import Document as DocxDocument
except ImportError:
    os.system("pip install python-docx")
    from docx import Document as DocxDocument
try:
    import openpyxl
except ImportError:
    os.system("pip install openpyxl")
    import openpyxl
try:
    import extract_msg
except ImportError:
    os.system("pip install extract-msg")
    import extract_msg
try:
    from pptx import Presentation
except ImportError:
    os.system("pip install python-pptx")
    from pptx import Presentation


# =====================================================
# CONFIGURATION
# =====================================================
FILTERED_DIR = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Filtrees"
OUTPUT_DIR = r"G:\Mon Drive\Projet SmarterPlan\Sales\Prospects\NCG\202512 Mission Déploiement IA interne\Résultats bruts\Archives_Extraites"
S3_BUCKET = "smarterplan-rag-prototype"
S3_TEXTRACT_PREFIX = "textract_temp/"
AWS_REGION = "eu-west-1"

# ── Paramètres de performance ──
MAX_CONCURRENT_UPLOADS = 20       # Threads pour upload S3
JOB_LAUNCH_RATE = 15              # Jobs lancés par seconde (quota augmenté)
POLL_INTERVAL = 3                 # Secondes entre chaque cycle de polling
POLL_BATCH_SIZE = 100             # Jobs vérifiés par cycle (quota augmenté)
MAX_RETRIES = 3                   # Retries sur throttling

# ── Checkpoint / reprise ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CHECKPOINT_FILE = os.path.join(SCRIPT_DIR, "extraction_checkpoint.json")

# ── AWS clients avec retry adaptatif ──
boto_config = Config(
    retries={"max_attempts": 10, "mode": "adaptive"},
    max_pool_connections=50
)
textract = boto3.client("textract", region_name=AWS_REGION, config=boto_config)
s3 = boto3.client("s3", region_name=AWS_REGION, config=boto_config)

# ── Logging ──
log = logging.getLogger(__name__)

def setup_logging():
    # Nettoie les anciens handlers si on relance dans le même interpréteur
    for handler in logging.root.handlers[:]:
        logging.root.removeHandler(handler)
        
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        handlers=[
            logging.StreamHandler(),
            logging.FileHandler("extraction.log", mode='w', encoding="utf-8")
        ]
    )

# ── Stats ──
stats = {
    "pdf_natif": 0, "pdf_ocr": 0, "word": 0, "excel": 0,
    "email": 0, "texte": 0, "pptx": 0, "image_ocr": 0,
    "erreurs": 0, "vides": 0, "skipped_checkpoint": 0
}


# =====================================================
# CHECKPOINT : reprise après interruption
# =====================================================
def load_checkpoint():
    completed = set()
    checkpoint_exists_and_valid = False
    
    if os.path.exists(CHECKPOINT_FILE):
        try:
            with open(CHECKPOINT_FILE, "r") as f:
                data = json.load(f).get("completed", [])
                if data:
                    completed.update(data)
                    checkpoint_exists_and_valid = True
        except:
            pass
            
    # Ne scanner le dossier QUE SI on n'a pas de checkpoint valide
    # Cela évite de scanner 10 000 fichiers à chaque lancement
    if not checkpoint_exists_and_valid and os.path.exists(OUTPUT_DIR):
        print("Recherche des fichiers déjà traités dans le dossier de sortie (scan initial)...")
        for root, _, files in os.walk(OUTPUT_DIR):
            for fname in files:
                if fname.endswith(".json"):
                    # Retrouver le chemin relatif original (sans .json)
                    rel_json = os.path.relpath(os.path.join(root, fname), OUTPUT_DIR)
                    completed.add(rel_json[:-5])
                    
    return completed

def save_checkpoint(completed_set):
    with open(CHECKPOINT_FILE, "w") as f:
        json.dump({"completed": list(completed_set)}, f)

completed_files = load_checkpoint()
if completed_files:
    log.info(f"♻️  Reprise : {len(completed_files)} fichiers déjà traités")


# =====================================================
# EXTRACTEURS DIRECTS (non-Textract)
# =====================================================
def extract_pdf_native(filepath):
    """Retourne (texte, is_native). is_native=False → nécessite Textract."""
    try:
        doc = fitz.open(filepath)
        text = ""
        page_count = doc.page_count
        for page in doc:
            text += page.get_text()
        doc.close()
        if len(text.strip()) > 50 * max(page_count, 1):
            return text.strip(), True
    except Exception:
        pass
    return "", False

def extract_docx(filepath):
    try:
        doc = DocxDocument(filepath)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception:
        try:
            with open(filepath, "rb") as f:
                raw = f.read()
            text = raw.decode("utf-8", errors="ignore")
            return "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        except Exception:
            return ""

def extract_excel(filepath):
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            texts.append(f"--- Feuille: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join([str(c) for c in row if c is not None])
                if row_text.strip():
                    texts.append(row_text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_email_msg(filepath):
    try:
        msg = extract_msg.Message(filepath)
        parts = []
        if msg.subject: parts.append(f"Objet: {msg.subject}")
        if msg.sender: parts.append(f"De: {msg.sender}")
        if msg.date: parts.append(f"Date: {msg.date}")
        if msg.body: parts.append(f"\n{msg.body}")
        return "\n".join(parts)
    except Exception:
        return ""

def extract_pptx(filepath):
    try:
        prs = Presentation(filepath)
        texts = []
        for i, slide in enumerate(prs.slides):
            texts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text_file(filepath):
    for enc in ["utf-8", "latin-1", "cp1252"]:
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except Exception:
            continue
    return ""


DIRECT_EXTRACTORS = {
    "word": extract_docx,
    "excel": extract_excel,
    "email": extract_email_msg,
    "texte": extract_text_file,
    "pptx": extract_pptx,
}

DIRECT_EXT_MAP = {
    ".doc": "word", ".docx": "word",
    ".xls": "excel", ".xlsx": "excel",
    ".msg": "email", ".eml": "email",
    ".txt": "texte", ".rtf": "texte", ".csv": "texte",
    ".ppt": "pptx", ".pptx": "pptx",
}

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp"}


# =====================================================
# HELPERS
# =====================================================
def save_extracted(rel_path, ext, text, output_dir):
    """Sauvegarde un résultat d'extraction en JSON."""
    if not text or len(text.strip()) < 20:
        stats["vides"] += 1
        return False

    output_data = {
        "source_file": rel_path,
        "source_extension": ext,
        "copropriete": rel_path.split(os.sep)[0],
        "dossier_parent": os.path.dirname(rel_path),
        "nom_fichier": os.path.basename(rel_path),
        "texte": text,
        "nb_caracteres": len(text)
    }
    output_path = os.path.join(output_dir, rel_path + ".json")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)
    return True


# =====================================================
# PHASE 0 : TRIAGE
# =====================================================
def triage_files():
    """Sépare fichiers en extraction directe vs OCR Textract."""
    direct_files = []   # (filepath, rel_path, ext, extractor_name)
    ocr_files = []      # (filepath, rel_path, ext)

    log.info("Phase 0: Triage des fichiers...")
    for root, dirs, filenames in os.walk(FILTERED_DIR):
        for fname in filenames:
            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, FILTERED_DIR)
            ext = Path(fname).suffix.lower()

            if rel_path in completed_files:
                stats["skipped_checkpoint"] += 1
                continue

            if ext in DIRECT_EXT_MAP:
                direct_files.append((filepath, rel_path, ext, DIRECT_EXT_MAP[ext]))
            elif ext == ".pdf":
                text, is_native = extract_pdf_native(filepath)
                if is_native:
                    # Stocker le texte déjà extrait pour éviter double lecture
                    direct_files.append((filepath, rel_path, ext, "pdf_natif"))
                else:
                    ocr_files.append((filepath, rel_path, ext))
            elif ext in IMAGE_EXTS:
                ocr_files.append((filepath, rel_path, ext))

    return direct_files, ocr_files


# =====================================================
# PHASE 1 : EXTRACTION DIRECTE
# =====================================================
def run_direct_extraction(direct_files):
    log.info(f"\nPhase 1: Extraction directe de {len(direct_files)} fichiers...")
    for filepath, rel_path, ext, etype in tqdm(direct_files, desc="Extraction directe"):
        if etype == "pdf_natif":
            text, _ = extract_pdf_native(filepath)
        else:
            text = DIRECT_EXTRACTORS[etype](filepath)

        if save_extracted(rel_path, ext, text, OUTPUT_DIR):
            stats[etype] += 1
            completed_files.add(rel_path)
        else:
            if not text or len(text.strip()) < 20:
                pass  # déjà compté dans save_extracted
            else:
                stats["erreurs"] += 1

    save_checkpoint(completed_files)
    log.info(f"✅ Phase 1 terminée")


# =====================================================
# PHASE 2 : UPLOAD S3 PARALLÈLE
# =====================================================
def upload_one(filepath, s3_key):
    try:
        s3.upload_file(filepath, S3_BUCKET, s3_key)
        return s3_key
    except Exception as e:
        log.error(f"Upload échoué {os.path.basename(filepath)}: {e}")
        return None

def batch_upload_s3(ocr_files):
    log.info(f"\nPhase 2: Upload S3 de {len(ocr_files)} fichiers ({MAX_CONCURRENT_UPLOADS} threads)...")
    s3_map = {}  # rel_path → (s3_key, ext)

    with ThreadPoolExecutor(max_workers=MAX_CONCURRENT_UPLOADS) as executor:
        futures = {}
        for filepath, rel_path, ext in ocr_files:
            s3_key = S3_TEXTRACT_PREFIX + rel_path.replace(os.sep, "/")
            fut = executor.submit(upload_one, filepath, s3_key)
            futures[fut] = (rel_path, s3_key, ext)

        for fut in tqdm(as_completed(futures), total=len(futures), desc="Upload S3"):
            rel_path, s3_key, ext = futures[fut]
            if fut.result():
                s3_map[rel_path] = (s3_key, ext)

    log.info(f"✅ {len(s3_map)}/{len(ocr_files)} fichiers sur S3")
    return s3_map


# =====================================================
# PHASE 3 : LANCEMENT RAFALE TEXTRACT
# =====================================================
def start_textract_job(s3_key, ext):
    """Lance un job. Retourne ('ASYNC', job_id) ou ('SYNC', text) ou None."""
    for attempt in range(MAX_RETRIES):
        try:
            if ext == ".pdf":
                resp = textract.start_document_text_detection(
                    DocumentLocation={"S3Object": {"Bucket": S3_BUCKET, "Name": s3_key}}
                )
                return ("ASYNC", resp["JobId"])
            else:
                # Images ≤10MB : API synchrone (plus rapide, pas de polling)
                obj = s3.get_object(Bucket=S3_BUCKET, Key=s3_key)
                img_bytes = obj["Body"].read()

                if len(img_bytes) > 10_000_000:
                    resp = textract.start_document_text_detection(
                        DocumentLocation={"S3Object": {"Bucket": S3_BUCKET, "Name": s3_key}}
                    )
                    return ("ASYNC", resp["JobId"])

                resp = textract.detect_document_text(Document={"Bytes": img_bytes})
                lines = [b["Text"] for b in resp.get("Blocks", []) if b["BlockType"] == "LINE"]
                return ("SYNC", "\n".join(lines))

        except textract.exceptions.ProvisionedThroughputExceededException:
            time.sleep(2 ** attempt)
        except Exception as e:
            log.error(f"Textract start error {s3_key}: {e}")
            return None
    return None

def launch_all_jobs(s3_map):
    log.info(f"\nPhase 3: Lancement de {len(s3_map)} jobs Textract (rate: {JOB_LAUNCH_RATE}/s)...")
    async_jobs = {}      # job_id → (rel_path, s3_key, ext)
    sync_results = {}    # rel_path → (text, ext)
    failed = []

    for i, (rel_path, (s3_key, ext)) in enumerate(
        tqdm(s3_map.items(), desc="Lancement Textract")
    ):
        result = start_textract_job(s3_key, ext)

        if result is None:
            failed.append(rel_path)
        elif result[0] == "ASYNC":
            async_jobs[result[1]] = (rel_path, s3_key, ext)
        elif result[0] == "SYNC":
            sync_results[rel_path] = (result[1], ext)

        # Rate limiting
        if (i + 1) % JOB_LAUNCH_RATE == 0:
            time.sleep(1.0)

    log.info(f"✅ {len(async_jobs)} async + {len(sync_results)} sync + {len(failed)} échecs")
    return async_jobs, sync_results, failed


# =====================================================
# PHASE 4 : POLLING DES RÉSULTATS ASYNC
# =====================================================
def collect_result(job_id):
    """Récupère toutes les pages de résultat d'un job terminé."""
    text_blocks = []
    kwargs = {"JobId": job_id}
    while True:
        result = textract.get_document_text_detection(**kwargs)
        for block in result.get("Blocks", []):
            if block["BlockType"] == "LINE":
                text_blocks.append(block["Text"])
        next_token = result.get("NextToken")
        if not next_token:
            break
        kwargs["NextToken"] = next_token
    return "\n".join(text_blocks)

def poll_all_jobs(async_jobs):
    log.info(f"\nPhase 4: Polling de {len(async_jobs)} jobs...")
    pending = dict(async_jobs)
    results = {}
    failed = []
    cycle = 0

    while pending:
        cycle += 1
        done_this_round = []
        job_ids = list(pending.keys())[:POLL_BATCH_SIZE]

        for job_id in job_ids:
            try:
                resp = textract.get_document_text_detection(JobId=job_id, MaxResults=1)
                status = resp["JobStatus"]

                if status == "SUCCEEDED":
                    rel_path, s3_key, ext = pending[job_id]
                    text = collect_result(job_id)
                    results[rel_path] = (text, ext)
                    done_this_round.append(job_id)
                elif status == "FAILED":
                    rel_path, _, _ = pending[job_id]
                    log.warning(f"FAILED: {rel_path}")
                    failed.append(rel_path)
                    done_this_round.append(job_id)

            except Exception as e:
                log.warning(f"Poll error: {e}")

        for jid in done_this_round:
            del pending[jid]

        if pending:
            done_n = len(results) + len(failed)
            total_n = done_n + len(pending)
            log.info(f"  Cycle {cycle}: {done_n}/{total_n} ({done_n/total_n*100:.0f}%) — {len(pending)} restants")
            time.sleep(POLL_INTERVAL)

    log.info(f"✅ Polling terminé: {len(results)} OK, {len(failed)} KO")
    return results, failed


# =====================================================
# PHASE 5 : NETTOYAGE S3
# =====================================================
def cleanup_s3(s3_map):
    log.info("\nPhase 5: Nettoyage S3...")
    keys = [s3_key for s3_key, _ in s3_map.values()]
    for i in range(0, len(keys), 1000):
        batch = keys[i:i+1000]
        try:
            s3.delete_objects(
                Bucket=S3_BUCKET,
                Delete={"Objects": [{"Key": k} for k in batch]}
            )
        except Exception as e:
            log.warning(f"Nettoyage S3 partiel: {e}")
    log.info(f"✅ {len(keys)} fichiers temp supprimés")


# =====================================================
# MAIN
# =====================================================
def main():
    # Initialisation du log (le mode 'w' s'occupe de vider le fichier proprement)
    setup_logging()

    print("=" * 60)
    print("  EXTRACTION OPTIMISÉE — PIPELINE PARALLÈLE")
    print("=" * 60)

    # ── Triage ──
    direct_files, ocr_files = triage_files()
    log.info(f"  {len(direct_files)} directs | {len(ocr_files)} OCR | {stats['skipped_checkpoint']} checkpoint")

    # ── Phase 1 : Extraction directe ──
    run_direct_extraction(direct_files)

    if not ocr_files:
        log.info("Aucun fichier OCR. Terminé !")
        print_report()
        return

    # ── Phase 2 : Upload S3 ──
    s3_map = batch_upload_s3(ocr_files)

    # ── Phase 3 : Lancement Textract ──
    async_jobs, sync_results, launch_failed = launch_all_jobs(s3_map)

    # Sauvegarder résultats synchrones immédiatement
    for rel_path, (text, ext) in sync_results.items():
        if save_extracted(rel_path, ext, text, OUTPUT_DIR):
            stats["image_ocr"] += 1
            completed_files.add(rel_path)
    save_checkpoint(completed_files)

    # ── Phase 4 : Polling async ──
    if async_jobs:
        ocr_results, poll_failed = poll_all_jobs(async_jobs)
        for rel_path, (text, ext) in ocr_results.items():
            if save_extracted(rel_path, ext, text, OUTPUT_DIR):
                if ext == ".pdf":
                    stats["pdf_ocr"] += 1
                else:
                    stats["image_ocr"] += 1
                completed_files.add(rel_path)
        stats["erreurs"] += len(poll_failed) + len(launch_failed)
        save_checkpoint(completed_files)

    # ── Phase 5 : Nettoyage ──
    cleanup_s3(s3_map)

    # ── Rapport ──
    print_report()

    if stats["erreurs"] == 0 and os.path.exists(CHECKPOINT_FILE):
        os.remove(CHECKPOINT_FILE)
    elif stats["erreurs"] > 0:
        ERROR_LOG = os.path.join(SCRIPT_DIR, "erreurs_extraction.txt")
        with open(ERROR_LOG, "w", encoding="utf-8") as f:
            f.write("LISTE DES FICHIERS EN ERREUR\n")
            f.write("="*30 + "\n")
            for rel_id in (launch_failed + (poll_failed if 'poll_failed' in locals() else [])):
                f.write(f"- {rel_id}\n")
        print(f"\n  ⚠️  {stats['erreurs']} erreurs. Liste détaillée dans : {ERROR_LOG}")
        print(f"  Relance le script pour retenter uniquement les fichiers échoués.")


def print_report():
    print("\n" + "=" * 60)
    print("  RAPPORT D'EXTRACTION")
    print("=" * 60)
    print(f"  PDF natifs (texte direct) : {stats['pdf_natif']}")
    print(f"  PDF scannés (Textract)    : {stats['pdf_ocr']}")
    print(f"  Word                      : {stats['word']}")
    print(f"  Excel                     : {stats['excel']}")
    print(f"  Emails                    : {stats['email']}")
    print(f"  Texte/CSV                 : {stats['texte']}")
    print(f"  PowerPoint                : {stats['pptx']}")
    print(f"  Images OCR (plans)        : {stats['image_ocr']}")
    print(f"  Fichiers vides/courts     : {stats['vides']}")
    print(f"  Erreurs                   : {stats['erreurs']}")
    t = stats['pdf_ocr'] + stats['image_ocr']
    print(f"\n  TOTAL Textract            : {t} fichiers")
    print(f"  Coût Textract estimé      : ~${t * 0.0015:.2f}")
    print(f"\n  📁 Résultats : {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
